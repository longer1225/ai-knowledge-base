import os
import tempfile
import json
from fastapi import UploadFile

from langchain_text_splitters import RecursiveCharacterTextSplitter
from backend.config import TEXT_SPLIT_CONFIG
from backend.core.exceptions import ParamException

# 🔥 已删除：DocumentChunk 导入（绝对不能在 Service 出现 Model）
from backend.utils.embedding_util import get_embedding
from backend.utils.logger import logger
from backend.utils.redis_cache import redis_cache

from backend.mapper.document_mapper import (
    insert_document,
    batch_insert_chunks,
    get_document_by_name
)

# ======================
# 分块器
# ======================
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=TEXT_SPLIT_CONFIG["chunk_size"],
    chunk_overlap=TEXT_SPLIT_CONFIG["chunk_overlap"],
    separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", "；", ";", "，", ","]
)

# ======================
# embedding
# ======================
embedding = get_embedding()

# ======================
# 支持格式
# ======================
ALLOWED_SUFFIX = [".txt", ".md", ".docx", ".pdf", ".csv", ".pptx"]

# ======================
# 🔥 Qwen 限制
# ======================
BATCH_SIZE = 25


def batch_embed(chunks):
    """手动分批 embedding"""
    all_vectors = []

    for i in range(0, len(chunks), BATCH_SIZE):
        batch = chunks[i:i + BATCH_SIZE]
        logger.debug(f"[Embedding] 处理 batch {i} ~ {i + len(batch)}")

        vectors = embedding.embed_documents(batch)
        all_vectors.extend(vectors)

    return all_vectors


def upload_document(file: UploadFile, user_id: int):
    logger.info(f"[DEBUG] 使用 embedding 模型类型: {type(embedding).__name__}")
    logger.info("[Service] 开始处理用户 %s 的文件：%s", user_id, file.filename)

    # ======================
    # 文件名去重
    # ======================
    exist_doc = get_document_by_name(user_id, file.filename)
    if exist_doc:
        doc_id = exist_doc["doc_id"]
        logger.info("[Service] 文件重复上传，直接返回 doc_id=%s", doc_id)
        return doc_id

    suffix = os.path.splitext(file.filename)[-1].lower()
    if suffix not in ALLOWED_SUFFIX:
        raise ParamException(f"仅支持：{', '.join(ALLOWED_SUFFIX)}")

    # ======================
    # 保存临时文件
    # ======================
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(file.file.read())
        tmp_path = tmp.name

    try:
        content = ""

        # ======================
        # 文件解析
        # ======================
        if suffix in [".txt", ".md"]:
            with open(tmp_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

        elif suffix == ".docx":
            from docx import Document
            doc_file = Document(tmp_path)
            content = "\n".join(
                [p.text for p in doc_file.paragraphs if p.text.strip()]
            )

        elif suffix == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(tmp_path)
            content = "\n".join([
                page.extract_text() or "" for page in reader.pages
            ])

        elif suffix == ".csv":
            import csv
            rows = []
            with open(tmp_path, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(row))
            content = "\n".join(rows)

        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(tmp_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            content = "\n".join(texts)

        # ======================
        # 内容校验
        # ======================
        if not content.strip():
            raise ParamException("文件内容为空或无法提取文本")

        # ======================
        # 分块
        # ======================
        chunks = text_splitter.split_text(content)

        if not chunks:
            raise ParamException("文本分块失败，内容过短或格式异常")

        logger.debug("[Service] 文件分块完成，共 %s 块", len(chunks))

        # ======================
        # 插入文档
        # ======================
        doc = insert_document(
            user_id=user_id,
            doc_name=file.filename,
            doc_type=suffix[1:],
            file_size=os.path.getsize(tmp_path)
        )

        doc_id = doc["doc_id"]

        # ======================
        # embedding（分批）
        # ======================
        vectors = batch_embed(chunks)

        if len(vectors) != len(chunks):
            raise Exception("embedding 结果数量与文本块不一致")

        # ======================
        # 🔥 修复：Service 只传字典，不创建 Model
        # ======================
        chunk_items = [
            {
                "doc_id": doc_id,
                "chunk_text": text,
                "chunk_embedding": vec,
                "chunk_index": idx
            }
            for idx, (text, vec) in enumerate(zip(chunks, vectors))
        ]

        # ======================
        # 入库
        # ======================
        chunk_dicts = batch_insert_chunks(chunk_items)

        # ======================
        # Redis缓存
        # ======================
        redis_cache.set(
            f"rag:doc:{doc_id}:chunks",
            json.dumps(chunk_dicts, ensure_ascii=False),
            ex=3600
        )

        # ======================
        # 全局版本号
        # ======================
        if not redis_cache.get("rag:version"):
            redis_cache.set("rag:version", 1)
        else:
            redis_cache.incr("rag:version")

        logger.info("[Service] 文档处理完成，doc_id=%s", doc_id)
        return doc_id

    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)